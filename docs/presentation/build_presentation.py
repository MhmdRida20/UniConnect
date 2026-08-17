"""
Builds the FYP defence deck: UniConnect_FYP_Presentation.pptx

Everything on the slides is generated from this file, so the deck can be
rebuilt after any edit rather than being patched by hand in PowerPoint. Run:

    python docs/presentation/build_presentation.py

Design notes
------------
The palette is lifted straight from wwwroot/css/site.css so the deck and the
product read as one thing. Typography is Segoe UI, which is the fallback the
site's own font stack names after Inter -- Inter is not installed on the
presentation machine, and a deck that silently substitutes a font is worse than
one that picks the safe member of the same stack.

Icons are the project's own Hugeicons, rendered out of Views/Shared/_Icons.cshtml
by render_icons.py, so the slides use the literal icons the portal uses.

Numbers on the evidence slide are the verified ones from docs/README.md. If the
codebase moves, re-derive them there first; this file is not the source of truth
for them.
"""
import os
import re
import subprocess
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import theme
from theme import animate, fill_alpha, set_transition, soft_shadow

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.abspath(os.path.join(HERE, "..", ".."))
ASSETS = os.path.join(HERE, "assets")
ICONS = os.path.join(ASSETS, "icons")
OUT = os.path.join(ROOT, "UniConnect_FYP_Presentation.pptx")

# ---------------------------------------------------------------- palette ---
INK        = "0f172a"
INK_SOFT   = "334155"
MUTED      = "64748b"
FAINT      = "94a3b8"
BORDER     = "e2e8f0"
BORDER_SFT = "eef2f6"
WHITE      = "ffffff"
BG         = "f6f9f7"

GREEN      = "16a34a"
GREEN_DK   = "15803d"
GREEN_DKR  = "166534"
GREEN_DEEP = "14532d"
GREEN_LT   = "22c55e"
GREEN_SOFT = "dcfce7"
GREEN_TINT = "f0fdf4"

TEAL       = "0d9488"
TEAL_SOFT  = "ccfbf1"
AMBER      = "d97706"
AMBER_SOFT = "fef3c7"
AMBER_TINT = "fffbeb"

FONT = "Segoe UI"
FONT_SB = "Segoe UI Semibold"

# ----------------------------------------------------------------- layout ---
SW, SH = 13.333, 7.5      # 16:9
M = 0.75                  # side margin
CW = SW - 2 * M           # content width
TOP = 1.72                # first content row
BOTTOM = 6.82             # last usable line
FOOT = 6.98


def rgb(h):
    return RGBColor.from_string(h.upper())


# ------------------------------------------------------------- primitives ---
def noshadow(shape):
    """python-pptx gives every autoshape a default drop shadow; the deck is flat."""
    shape.shadow.inherit = False
    return shape


def since(slide, n):
    """Every shape added to `slide` after index `n` — one animation wave."""
    return list(slide.shapes)[n:]


def waves(slide, marks):
    """
    Turn shape-count marks into animation waves.

    Each slide records len(slide.shapes) at the points where one visual group
    ends and the next begins; this slices the shape list on those boundaries, so
    a group animates together no matter how many shapes it is built from.
    """
    shapes = list(slide.shapes)
    bounds = list(marks) + [len(shapes)]
    return [shapes[a:b] for a, b in zip(bounds, bounds[1:]) if b > a]


def rect(slide, x, y, w, h, fill=None, line=None, lw=1.0, radius=None,
         shadow=False, alpha=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    noshadow(s)
    if radius:
        # adjustment is a fraction of half the shorter side
        s.adjustments[0] = min(0.5, radius / (min(w, h) / 2.0) * 0.5)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = rgb(fill)
        if alpha is not None:
            fill_alpha(s, alpha)
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = rgb(line)
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.text_frame.text = ""
    if shadow:
        soft_shadow(s)
    return s


def line_h(slide, x, y, w, colour=BORDER, weight=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Pt(weight))
    noshadow(s)
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(colour)
    s.line.fill.background()
    return s


def tbox(slide, x, y, w, h, anchor="top"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP,
                          "mid": MSO_ANCHOR.MIDDLE,
                          "bot": MSO_ANCHOR.BOTTOM}[anchor]
    return tf


def para(tf, text, size=12, colour=INK_SOFT, bold=False, font=FONT, first=False,
         before=0, after=0, align=None, spacing=None, tracking=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align:
        p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER,
                       "r": PP_ALIGN.RIGHT}[align]
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    if spacing:
        p.line_spacing = spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = rgb(colour)
    if tracking:
        # letter-spacing has no python-pptx API; spc is hundredths of a point
        r.font._rPr.set("spc", str(int(tracking * 100)))
    return p


def icon(slide, name, colour, x, y, size):
    path = os.path.join(ICONS, f"{name}_{colour}.png")
    if not os.path.exists(path):
        raise FileNotFoundError(path)
    return slide.shapes.add_picture(path, Inches(x), Inches(y),
                                    Inches(size), Inches(size))


def picture(slide, path, x, y, w=None, h=None):
    kw = {}
    if w:
        kw["width"] = Inches(w)
    if h:
        kw["height"] = Inches(h)
    return slide.shapes.add_picture(path, Inches(x), Inches(y), **kw)


# ----------------------------------------------------------------- chrome ---
def new_slide(prs, dark=False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = os.path.join(ASSETS, "bg_dark.png" if dark else "bg_light.png")
    picture(s, bg, 0, 0, w=SW, h=SH)
    return s


def heading(slide, eyebrow, title, sub=None, dark=False):
    # Titles are kept to one line on purpose: a wrapped title pushes into the
    # accent rule and the first content row, and every slide sharing one grid
    # is what makes the deck read as designed rather than assembled.
    if len(title) > 60:
        raise ValueError(f"title wraps at 30pt ({len(title)} chars): {title!r}")

    tf = tbox(slide, M, 0.52, CW, 0.28)
    para(tf, eyebrow.upper(), 10.5, GREEN_LT if dark else GREEN, bold=True,
         font=FONT_SB, first=True, tracking=1.6)

    ty = 0.84
    tf = tbox(slide, M, ty, CW, 0.66)
    para(tf, title, 30, WHITE if dark else INK, bold=True, font=FONT_SB,
         first=True, spacing=1.0)

    y = 1.52
    if sub:
        tf = tbox(slide, M, y, CW - 1.0, 0.3)
        para(tf, sub, 12.5, GREEN_SOFT if dark else MUTED, first=True)
        y += 0.36
    line_h(slide, M, y, 1.5, GREEN if not dark else GREEN_LT, 2.4)
    return y + 0.30


def flag_pill(slide, text, icon_name="i-check-badge", w=3.75):
    """A highlight pill on the title row, for a claim worth showing, not burying."""
    x, y, h = SW - M - w, 0.90, 0.52
    rect(slide, x, y, w, h, fill=GREEN_TINT, line=GREEN_SOFT, radius=0.14)
    icon(slide, icon_name, "16a34a", x + 0.22, y + (h - 0.26) / 2, 0.26)
    tf = tbox(slide, x + 0.60, y + 0.06, w - 0.80, h - 0.12, anchor="mid")
    para(tf, text, 10.5, GREEN_DEEP, bold=True, font=FONT_SB, first=True, spacing=1.08)


def footer(slide, n, dark=False):
    tf = tbox(slide, M, FOOT, CW * 0.7, 0.24)
    para(tf, "UniConnect  ·  Final Year Project  ·  University of Sciences and Arts in Lebanon",
         8.5, FAINT if not dark else "6b8f7a", first=True)
    tf = tbox(slide, SW - M - 1.0, FOOT, 1.0, 0.24)
    para(tf, str(n), 8.5, FAINT if not dark else "6b8f7a", first=True, align="r")


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


# ------------------------------------------------------------- components ---
def card(slide, x, y, w, h, title, body, icon_name=None, icon_colour=GREEN,
         fill=WHITE, border=BORDER, title_size=13.5, body_size=10.5,
         title_colour=INK, body_colour=MUTED):
    rect(slide, x, y, w, h, fill=fill, line=border, radius=0.12, shadow=True)
    pad = 0.26
    ty = y + pad
    if icon_name:
        chip = rect(slide, x + pad, ty, 0.44, 0.44, fill=GREEN_SOFT, radius=0.1)
        icon(slide, icon_name, icon_colour, x + pad + 0.10, ty + 0.10, 0.24)
        tx, tw = x + pad + 0.60, w - 2 * pad - 0.60
        tf = tbox(slide, tx, ty + 0.04, tw, 0.36, anchor="mid")
        para(tf, title, title_size, title_colour, bold=True, font=FONT_SB, first=True)
        ty += 0.60
    else:
        tf = tbox(slide, x + pad, ty, w - 2 * pad, 0.30)
        para(tf, title, title_size, title_colour, bold=True, font=FONT_SB, first=True)
        ty += 0.36
    tf = tbox(slide, x + pad, ty, w - 2 * pad, h - (ty - y) - pad)
    para(tf, body, body_size, body_colour, first=True, spacing=1.20)


def bullet(tf, text, size=11.5, colour=INK_SOFT, first=False, before=6, bold_lead=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(before)
    p.space_after = Pt(0)
    p.line_spacing = 1.18
    dot = p.add_run()
    dot.text = "—  "
    dot.font.size = Pt(size)
    dot.font.name = FONT
    dot.font.color.rgb = rgb(GREEN)
    dot.font.bold = True
    if bold_lead:
        r = p.add_run()
        r.text = bold_lead
        r.font.size = Pt(size)
        r.font.name = FONT_SB
        r.font.bold = True
        r.font.color.rgb = rgb(INK)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.name = FONT
    r.font.color.rgb = rgb(colour)
    return p


def callout(slide, x, y, w, h, text, icon_name="i-idea", tone="green"):
    if tone == "green":
        fill, brd, ic, tc = GREEN_TINT, GREEN_SOFT, GREEN, GREEN_DEEP
    elif tone == "amber":
        fill, brd, ic, tc = AMBER_TINT, "fde68a", AMBER, "78350f"
    else:
        fill, brd, ic, tc = BG, BORDER, MUTED, INK_SOFT
    rect(slide, x, y, w, h, fill=fill, line=brd, radius=0.10)
    icon(slide, icon_name, ic.lstrip("#"), x + 0.24, y + (h - 0.26) / 2, 0.26)
    tf = tbox(slide, x + 0.66, y + 0.14, w - 0.92, h - 0.28, anchor="mid")
    para(tf, text, 11, tc, first=True, spacing=1.16)


def stat(slide, x, y, w, h, value, label):
    rect(slide, x, y, w, h, fill=WHITE, line=BORDER, radius=0.12, shadow=True)
    # a short accent rule under the number, so the tile reads as a designed unit
    rect(slide, x + w / 2 - 0.28, y + h - 0.60, 0.56, 0.035, fill=GREEN_SOFT)
    tf = tbox(slide, x + 0.20, y + 0.18, w - 0.40, 0.60, anchor="mid")
    para(tf, value, 34, GREEN_DEEP, bold=True, font=FONT_SB, first=True, align="c")
    tf = tbox(slide, x + 0.14, y + h - 0.46, w - 0.28, 0.40, anchor="top")
    para(tf, label, 9.5, MUTED, first=True, align="c", spacing=1.1)


# ------------------------------------------------------------ backgrounds ---
def make_backgrounds():
    theme.make_backgrounds(ASSETS)


# ================================================================= slides ===
def slide_title(prs):
    s = new_slide(prs, dark=True)
    COL = 6.55            # the text column; the product shot takes the rest

    # ---- product composition, right ---------------------------------------
    # A real screenshot of both clients does more for a mixed audience than any
    # amount of type. The phone overlaps the browser so the two read as one
    # product rather than two unrelated pictures.
    n_shot = len(s.shapes)
    web = os.path.join(ASSETS, "web_home_framed.png")
    phone = os.path.join(ASSETS, "mobile_home_framed.png")
    if os.path.exists(web):
        pic = picture(s, web, 7.50, 1.74, w=5.40)
        soft_shadow(pic, blur=30, dist=10, alpha=0.34)
    if os.path.exists(phone):
        pic = picture(s, phone, 6.86, 3.02, h=3.66)
        soft_shadow(pic, blur=26, dist=9, alpha=0.38)

    # ---- brand mark -------------------------------------------------------
    n0 = len(s.shapes)
    glow = os.path.join(ASSETS, "glow.png")
    if os.path.exists(glow):
        picture(s, glow, M - 0.92, 0.72, w=2.70, h=2.70)
    rect(s, M, 1.44, 0.88, 0.88, fill=GREEN, radius=0.23)
    icon(s, "i-graduation", "ffffff", M + 0.21, 1.65, 0.46)
    n1 = len(s.shapes)

    tf = tbox(s, M, 2.62, COL, 0.30)
    para(tf, "FINAL YEAR PROJECT  ·  BACHELOR OF COMPUTER SCIENCE", 10.5,
         GREEN_LT, bold=True, font=FONT_SB, first=True, tracking=1.6)

    tf = tbox(s, M, 2.94, COL, 1.10)
    para(tf, "UniConnect", 62, WHITE, bold=True, font=FONT_SB, first=True)
    n2 = len(s.shapes)

    tf = tbox(s, M, 4.06, COL, 0.78)
    para(tf, "An Integration-First Multi-University\nDigital Services Platform",
         18.5, GREEN_SOFT, first=True, spacing=1.18)

    line_h(s, M, 4.98, 1.9, GREEN_LT, 2.6)
    n3 = len(s.shapes)

    tf = tbox(s, M, 5.24, COL, 0.72)
    para(tf, "Mohamad Sabbagh  ·  1123206", 13, WHITE, bold=True,
         font=FONT_SB, first=True)
    para(tf, "Mohamad Ali Rida  ·  1122824", 13, WHITE, bold=True,
         font=FONT_SB, before=4)

    tf = tbox(s, M, 6.20, COL, 0.86)
    para(tf, "Supervisor:  Zahraa Sleem", 11.5, GREEN_SOFT, first=True)
    para(tf, "University of Sciences and Arts in Lebanon  ·  "
             "Faculty of Sciences and Arts", 10.5, "8fd0a8", before=4)
    para(tf, "18 August 2026", 10.5, "8fd0a8", before=3)

    set_transition(s, "fade")
    animate(s, waves(s, [n0, n1, n2, n3]) + [since(s, n_shot)[:n0 - n_shot]],
            dur=520, gap=190)

    notes(s, """
SPEAKER A  ·  20 sec  ·  45 words

Good afternoon. I'm Mohamad Ali Rida, and this is Mohamad Sabbagh.

You've just seen UniConnect running. In the next ten minutes we'll cover the
three things a demo can't show you: why we built it this way, what the
architecture actually guarantees, and where it falls short.

(Don't read the title slide aloud. Move on.)
""")
    return s


def slide_problem(prs):
    s = new_slide(prs)
    y = heading(s, "The problem",
                "Universities run the records. Students improvise the rest.")

    ch = 3.55
    lw = (CW - 0.34) / 2
    m0 = len(s.shapes)

    # left: what the institution already owns
    rect(s, M, y, lw, ch, fill=WHITE, line=BORDER, radius=0.14, shadow=True)
    rect(s, M, y, lw, 0.06, fill=GREEN, radius=None)
    tf = tbox(s, M + 0.32, y + 0.30, lw - 0.64, 0.3)
    para(tf, "What the SIS and LMS already own", 14, INK, bold=True,
         font=FONT_SB, first=True)
    tf = tbox(s, M + 0.32, y + 0.72, lw - 0.64, 2.5)
    for i, t in enumerate(["Student identity and university records",
                           "Course catalogue and enrolment",
                           "Coursework, submissions and grades"]):
        bullet(tf, t, 12, INK_SOFT, first=(i == 0), before=0 if i == 0 else 9)
    tf = tbox(s, M + 0.32, y + ch - 0.86, lw - 0.64, 0.68)
    para(tf, "Conservative by design, and rightly so. This is the institution's "
             "system of record, and no university should let a student project "
             "write to it.", 10.5, MUTED, first=True, spacing=1.16)

    # right: what students improvise
    m1 = len(s.shapes)
    x2 = M + lw + 0.34
    rect(s, x2, y, lw, ch, fill=WHITE, line=BORDER, radius=0.14, shadow=True)
    rect(s, x2, y, lw, 0.06, fill=AMBER, radius=None)
    tf = tbox(s, x2 + 0.32, y + 0.30, lw - 0.64, 0.3)
    para(tf, "What students are left to improvise", 14, INK, bold=True,
         font=FONT_SB, first=True)
    tf = tbox(s, x2 + 0.32, y + 0.72, lw - 0.64, 2.5)
    for i, t in enumerate(["Study partners, in a group chat anyone can join",
                           "Lifts to campus, arranged person to person",
                           "Attendance, by paper roll-call or a shared code",
                           "Internships, on job boards that never saw a transcript"]):
        bullet(tf, t, 12, INK_SOFT, first=(i == 0), before=0 if i == 0 else 9)
    tf = tbox(s, x2 + 0.32, y + ch - 0.86, lw - 0.64, 0.68)
    para(tf, "Convenient, and completely unaware of the university. That is the "
             "gap this project addresses.", 10.5, MUTED, first=True, spacing=1.16)

    m2 = len(s.shapes)
    by = y + ch + 0.26
    rect(s, M, by, CW, 0.74, fill=GREEN_DEEP, radius=0.12, shadow=True)
    icon(s, "i-alert", "d97706", M + 0.30, by + 0.24, 0.28)
    tf = tbox(s, M + 0.76, by + 0.12, CW - 1.1, 0.5, anchor="mid")
    para(tf, "None of these tools know who is actually enrolled — so none of them can enforce "
             "“only classmates in this course” or “only an enrolled student can be marked present.”",
         12.5, WHITE, first=True, spacing=1.14)

    footer(s, 2)
    set_transition(s, "fade")
    animate(s, waves(s, [m0, m1, m2]))
    notes(s, """
SPEAKER A  ·  50 sec  ·  120 words

Every university runs an SIS and an LMS. They own identity, enrolment and
grades, and they're deliberately conservative — no institution lets a student
project write into its record system.

Now look at the right column. Everything around campus life happens somewhere
else: a WhatsApp group for a study partner, a message for a lift, paper for
attendance, a public job board for an internship.

Those tools work. The problem is the bottom line — none of them know who is
enrolled. A WhatsApp group can't tell a classmate from a stranger with the link.
A shared attendance code can't tell a student in the room from one who was sent
a screenshot.

That's the gap. Not a missing app — a missing guarantee.
""")
    return s


def slide_solution(prs):
    s = new_slide(prs)
    y = heading(s, "The solution",
                "Nine services. One platform. No academic data owned.",
                sub="Each university enables or disables the six student-facing "
                    "services independently.")

    services = [
        ("i-qr", "Smart Attendance", "Time-bound QR sessions, verified on five checks."),
        ("i-users", "Study Groups", "Course-scoped groups with real-time chat."),
        ("i-car", "Ride Sharing", "Campus-only lifts with live GPS tracking."),
        ("i-flag", "Clubs & Organizations", "Membership roles, announcements, events with RSVP."),
        ("i-briefcase", "Internship Matching", "Scored against completed coursework."),
        ("i-ticket", "Complaints & Ticketing", "Routed to the right department, tracked to closure."),
    ]
    cw = (CW - 2 * 0.28) / 3
    ch = 1.42
    marks = []
    for i, (ic, t, b) in enumerate(services):
        if i % 3 == 0:
            marks.append(len(s.shapes))
        cx = M + (i % 3) * (cw + 0.28)
        cy = y + (i // 3) * (ch + 0.26)
        card(s, cx, cy, cw, ch, t, b, icon_name=ic, title_size=13, body_size=10.5)

    marks.append(len(s.shapes))
    by = y + 2 * ch + 0.26 + 0.30
    rect(s, M, by, CW, 0.92, fill=WHITE, line=BORDER, radius=0.12, shadow=True)
    tf = tbox(s, M + 0.32, by + 0.16, 2.7, 0.6, anchor="mid")
    para(tf, "Shared by every module", 11.5, INK, bold=True, font=FONT_SB, first=True)
    para(tf, "always on, not optional", 9.5, MUTED, before=2)
    px, step = M + 3.30, 2.80
    for ic, name in [("i-bell", "Notifications"), ("i-chart", "Reporting"),
                     ("i-file", "Audit Logging")]:
        icon(s, ic, "64748b", px, by + 0.32, 0.28)
        tf = tbox(s, px + 0.40, by + 0.28, step - 0.5, 0.36, anchor="mid")
        para(tf, name, 12, INK_SOFT, bold=True, font=FONT_SB, first=True)
        px += step

    footer(s, 3)
    set_transition(s, "fade")
    animate(s, waves(s, marks))
    notes(s, """
SPEAKER A  ·  45 sec  ·  90 words

UniConnect is one platform carrying nine services.

Six are student-facing — you saw most of them in the demo. Every university
enables or disables these six independently; one that only wants attendance
turns the other five off.

Three are cross-cutting: notifications, reporting, audit logging. Shared
infrastructure — which is why adding a service doesn't mean rebuilding a
notification system for the fourth time.

The important phrase is the last one in the title: no academic data owned.
That's the next slide.
""")
    return s


def slide_architecture(prs):
    s = new_slide(prs)
    y = heading(s, "Architecture",
                "One interface is the only thing the platform depends on.")

    # clients
    mA = len(s.shapes)
    cw = 3.5
    for i, (t, sub) in enumerate([("Web Portal", "Razor views  ·  cookie authentication"),
                                  ("Mobile App", ".NET MAUI  ·  JWT bearer")]):
        cx = M + 1.05 + i * (cw + 0.5)
        rect(s, cx, y, cw, 0.72, fill=WHITE, line=BORDER, radius=0.10, shadow=True)
        tf = tbox(s, cx + 0.22, y + 0.10, cw - 0.44, 0.52, anchor="mid")
        para(tf, t, 12.5, INK, bold=True, font=FONT_SB, first=True)
        para(tf, sub, 9.5, MUTED, before=1)

    ay = y + 0.80
    tf = tbox(s, M + 1.05, ay, 8.5, 0.22)
    para(tf, "▼", 10, FAINT, first=True, align="c")

    # backend
    mB = len(s.shapes)
    by = y + 1.02
    rect(s, M, by, CW, 0.96, fill=GREEN, radius=0.12, shadow=True)
    tf = tbox(s, M + 0.32, by + 0.14, 6.4, 0.68, anchor="mid")
    para(tf, "ASP.NET Core 8 backend", 15, WHITE, bold=True, font=FONT_SB, first=True)
    para(tf, "Core Platform  ·  9 service modules  ·  6 SignalR hubs  ·  6 background services",
         10.5, GREEN_SOFT, before=3)
    dbx = SW - M - 3.5
    rect(s, dbx, by + 0.16, 3.2, 0.64, fill=GREEN_DEEP, radius=0.10)
    tf = tbox(s, dbx + 0.18, by + 0.22, 2.84, 0.52, anchor="mid")
    para(tf, "UniConnect database", 11, WHITE, bold=True, font=FONT_SB, first=True)
    para(tf, "37 tables  ·  the only thing it writes to", 9, "a7e8c0", before=1)

    # adapter
    mC = len(s.shapes)
    dy = by + 1.08
    tf = tbox(s, M, dy, CW, 0.24)
    para(tf, "▼   reads academic data through   ▼", 9.5, MUTED, first=True, align="c")

    ady = dy + 0.26
    rect(s, M + 1.6, ady, CW - 3.2, 0.86, fill=WHITE, line=TEAL, lw=2.0,
         radius=0.12, shadow=True)
    icon(s, "i-cloud-sync", "16a34a", M + 2.0, ady + 0.24, 0.38)
    tf = tbox(s, M + 2.56, ady + 0.14, CW - 4.4, 0.6, anchor="mid")
    para(tf, "IUniversityProvider", 15.5, GREEN_DEEP, bold=True, font=FONT_SB, first=True)
    para(tf, "one adapter interface  ·  resolved per university at runtime", 10, MUTED, before=2)

    # universities
    mD = len(s.shapes)
    uy = ady + 1.00
    uw = 4.2
    for i, (t, sub) in enumerate([("Demo University API", "the project's own simulated tenant"),
                                  ("Partner University API", "a real, live production system")]):
        ux = M + 1.15 + i * (uw + 0.7)
        rect(s, ux, uy, uw, 0.70, fill=WHITE, line=BORDER, radius=0.10, shadow=True)
        tf = tbox(s, ux + 0.22, uy + 0.09, uw - 1.35, 0.52, anchor="mid")
        para(tf, t, 12, INK, bold=True, font=FONT_SB, first=True)
        para(tf, sub, 9.5, MUTED, before=1)
        rect(s, ux + uw - 1.10, uy + 0.20, 0.92, 0.30, fill=GREEN_SOFT, radius=0.08)
        tf = tbox(s, ux + uw - 1.10, uy + 0.22, 0.92, 0.26, anchor="mid")
        para(tf, "READ ONLY", 7.5, GREEN_DEEP, bold=True, font=FONT_SB,
             first=True, align="c", tracking=0.5)

    mE = len(s.shapes)
    callout(s, M, uy + 0.78, CW, 0.62,
            "Onboarding a new university means implementing one interface. "
            "Not one of the nine services changes.", "i-check-badge", "green")

    footer(s, 4)
    set_transition(s, "fade")
    # built bottom-up, because that is the order the speaker reads it in
    w = waves(s, [mA, mB, mC, mD, mE])
    animate(s, [w[3], w[2], w[1], w[0], w[4]], dur=460, gap=210)
    notes(s, """
SPEAKER A  ·  75 sec  ·  170 words  ·  THE KEY SLIDE — do not rush it

Read this diagram from the bottom.

At the bottom, university APIs. We never write to them. We read identity,
courses and enrolment; the university stays the system of record.

Above them, the single idea this project rests on: one interface,
IUniversityProvider. Every service that needs academic data asks that interface.
No service anywhere knows which university it's talking to, or what shape that
API has.

Above that, the backend — core platform, nine modules, SignalR hubs. It writes
to one database: ours, thirty-seven tables, holding only what we must own.
No grades. No transcripts.

At the top, both clients hit that same backend.

Now the payoff. We didn't just design this and hope. We implemented the
interface twice — once for our own simulated university, once for a real partner
university whose API is genuinely different: different field names, course data
nested inside sections, no bulk student endpoint.

Supporting it took one new class. Not one of the nine services changed. That's
the claim this project is really making.
""")
    return s


def slide_stack(prs):
    s = new_slide(prs)
    y = heading(s, "Technology", "Two clients, one backend, one identity.")

    items = [
        ("i-dashboard", "Backend", "ASP.NET Core 8 MVC in C#. Razor views for the web portal and a JSON API for mobile, sharing one service layer so the two clients cannot disagree about a rule."),
        ("i-book", "Data", "Entity Framework Core over SQL Server. 37 tables built by 33 migrations, with unique indexes and row versions enforcing the rules that matter."),
        ("i-bell", "Real-time", "SignalR across 6 hubs: group chat, notifications, live ride tracking, attendance, clubs and tickets."),
        ("i-id", "Identity", "ASP.NET Core Identity. Cookie authentication for the web session, JWT bearer for mobile — one account per student, either way in."),
        ("i-target", "Matching", "TF-IDF and cosine similarity over skills, completed coursework, major and interests."),
        ("i-check-circle", "Testing", "xUnit against a real file-backed SQLite database — 290 tests covering the highest-risk rules."),
    ]
    cw = (CW - 2 * 0.28) / 3
    ch = 1.72
    marks = []
    for i, (ic, t, b) in enumerate(items):
        if i % 3 == 0:
            marks.append(len(s.shapes))
        cx = M + (i % 3) * (cw + 0.28)
        cy = y + (i // 3) * (ch + 0.26)
        card(s, cx, cy, cw, ch, t, b, icon_name=ic, title_size=13, body_size=10)

    marks.append(len(s.shapes))
    by = y + 2 * ch + 0.26 + 0.22
    rect(s, M, by, CW, 0.56, fill=GREEN_TINT, line=GREEN_SOFT, radius=0.10)
    tf = tbox(s, M + 0.30, by + 0.08, CW - 0.6, 0.40, anchor="mid")
    para(tf, "Tooling:   Visual Studio 2022    ·    Git and GitHub    ·    EF Core migrations    ·    "
             "ZXing for QR scanning    ·    SMTP notifications    ·    .NET MAUI for Android and Windows",
         10.5, GREEN_DEEP, first=True)

    footer(s, 5)
    set_transition(s, "fade")
    animate(s, waves(s, marks))
    notes(s, """
SPEAKER A  ·  40 sec  ·  95 words  ·  Move briskly — reference slide

The stack is deliberately boring, and that's a decision.

ASP.NET Core 8 in C#. The web portal is server-rendered Razor; mobile talks to a
JSON API. Both go through the same service layer — so web and mobile can never
silently disagree about whether a student may join a group.

EF Core over SQL Server. SignalR, six hubs.

Identity is worth a sentence: cookies for web, JWT for mobile, both from the
same Identity store. One student, one account.

And the tests run against a real database — I'll come back to that.

[HAND OVER TO SPEAKER B]
""")
    return s


def slide_attendance(prs):
    s = new_slide(prs)
    y = heading(s, "Feature 01  ·  Smart Attendance",
                "Five checks, not one scanned code.")

    steps = [
        ("i-clock", "Token still valid", "The session's QR expires on a timer set by the instructor."),
        ("i-check-badge", "Actually enrolled", "Verified live through the adapter, against the university's own record."),
        ("i-copy", "Not already checked in", "A unique index on (session, student) makes a duplicate impossible."),
        ("i-location", "Inside the room", "Haversine distance against the session's GPS radius."),
        ("i-id", "Device not reused", "The same device cannot check in twice for one session."),
    ]
    shot_band = 2.55
    sw_ = (CW - shot_band - 4 * 0.20) / 5
    step_h = 2.62
    marks = []
    for i, (ic, t, b) in enumerate(steps):
        marks.append(len(s.shapes))
        cx = M + i * (sw_ + 0.20)
        rect(s, cx, y, sw_, step_h, fill=WHITE, line=BORDER, radius=0.12, shadow=True)
        rect(s, cx, y, sw_, 0.05, fill=GREEN)
        rect(s, cx + 0.20, y + 0.26, 0.34, 0.34, fill=GREEN_DEEP, radius=0.09)
        tf = tbox(s, cx + 0.20, y + 0.28, 0.34, 0.30, anchor="mid")
        para(tf, str(i + 1), 11, WHITE, bold=True, font=FONT_SB, first=True, align="c")
        icon(s, ic, "16a34a", cx + sw_ - 0.54, y + 0.26, 0.30)
        tf = tbox(s, cx + 0.20, y + 0.74, sw_ - 0.40, 0.46)
        para(tf, t, 11.5, INK, bold=True, font=FONT_SB, first=True, spacing=1.08)
        tf = tbox(s, cx + 0.20, y + 1.30, sw_ - 0.40, 1.10)
        para(tf, b, 9.5, MUTED, first=True, spacing=1.16)

    marks.append(len(s.shapes))
    left_w = CW - shot_band - 0.20
    ry = y + step_h + 0.22
    rect(s, M, ry, left_w, 0.62, fill=GREEN_DEEP, radius=0.10, shadow=True)
    tf = tbox(s, M + 0.28, ry + 0.10, left_w - 0.5, 0.42, anchor="mid")
    para(tf, "All five pass  →  recorded as Present, or Late outside the grace window",
         11.5, WHITE, bold=True, font=FONT_SB, first=True)

    mCall = len(s.shapes)
    callout(s, M, ry + 0.78, left_w, 0.86,
            "A card or code-only system checks one factor. Forward the code to a friend "
            "and it passes. Here it fails at step 4, and again at step 5.",
            "i-alert", "amber")

    shot = os.path.join(ASSETS, "mobile_checkin_framed.png")
    if os.path.exists(shot):
        from PIL import Image
        with Image.open(shot) as im:
            ratio = im.width / im.height
        sh = 4.24
        sw_shot = sh * ratio
        pic = picture(s, shot, SW - M - shot_band + (shot_band - sw_shot) / 2, y, h=sh)
        soft_shadow(pic, blur=22, dist=7, alpha=0.22)

    footer(s, 6)
    set_transition(s, "fade")
    # one wave per check: the build is the argument
    animate(s, waves(s, marks + [mCall]), dur=340, gap=150)
    notes(s, """
SPEAKER B  ·  65 sec  ·  150 words

This is the module we went deepest on.

The standard approach to digital attendance is a code or a card. One factor. And
one factor has an obvious failure: the student photographs the code and sends it
to a friend at home.

We check five things. Is the token still valid — sessions expire on a timer. Is
the student actually enrolled — that goes live through the adapter to the
university's own record. Have they already checked in — that's a unique index in
the database, so it's impossible by construction, not just code we could forget
to call. Are they inside the classroom radius. And has this device already
checked someone in for this session.

Pass all five: Present, or Late outside the grace window.

So back to the screenshot attack. The friend at home has a valid token and is
genuinely enrolled. They fail on GPS. Borrow a classmate's phone in the room,
and they fail on device reuse.
""")
    return s


def slide_internships(prs):
    s = new_slide(prs)
    y = heading(s, "Feature 02  ·  Internship Matching",
                "Matched on coursework, not on claims.")
    # CONFIRM THE WORDING BEFORE PRESENTING -- state only what is actually agreed.
    flag_pill(s, "Being adopted by our university\nfor internship placement")

    m0 = len(s.shapes)
    lw, ch = 6.30, 3.60
    rect(s, M, y, lw, ch, fill=WHITE, line=BORDER, radius=0.14, shadow=True)
    tf = tbox(s, M + 0.32, y + 0.26, lw - 0.64, 0.3)
    para(tf, "How the 0–100 score is composed", 13, INK, bold=True,
         font=FONT_SB, first=True)

    # One series, one colour: these are five parts of a single weighting, not
    # five entities needing their own identity, so a ramp would double-encode
    # bar length as hue and buy nothing.
    weights = [("Skills", 35), ("Completed courses", 25), ("Major", 20),
               ("Career interests", 10), ("Location", 10)]
    bar_x, bar_max_w, top = M + 2.05, 3.15, y + 0.82

    # Drawn in three passes so the bars can animate on their own: the empty
    # tracks and labels arrive first, the bars then sweep out from the baseline,
    # and the numbers land last. Same z-order as before -- fills are still added
    # after their tracks, so they sit on top.
    for i, (name, _) in enumerate(weights):
        ry = top + i * 0.44
        tf = tbox(s, M + 0.32, ry - 0.02, 1.65, 0.28, anchor="mid")
        para(tf, name, 10.5, INK_SOFT, first=True)
        rect(s, bar_x, ry + 0.02, bar_max_w, 0.20, fill=BORDER_SFT, radius=0.05)

    m_bars = len(s.shapes)
    for i, (_, val) in enumerate(weights):
        ry = top + i * 0.44
        rect(s, bar_x, ry + 0.02, bar_max_w * val / 35.0, 0.20, fill=GREEN, radius=0.05)

    m_vals = len(s.shapes)
    for i, (_, val) in enumerate(weights):
        ry = top + i * 0.44
        tf = tbox(s, bar_x + bar_max_w + 0.14, ry - 0.02, 0.55, 0.28, anchor="mid")
        para(tf, str(val), 11.5, GREEN_DEEP, bold=True, font=FONT_SB, first=True)

    tf = tbox(s, M + 0.32, y + ch - 0.66, lw - 0.64, 0.52)
    para(tf, "Weights sum to 100. Skills, courses, major and interests are each "
             "compared using TF-IDF and cosine similarity.", 10, MUTED,
         first=True, spacing=1.16)

    m1 = len(s.shapes)
    rx = M + lw + 0.30
    rwid = CW - lw - 0.30
    rh = (ch - 0.16) / 2
    card(s, rx, y, rwid, rh,
         "Read live, through the adapter",
         "The course component is the student's actual completed coursework, "
         "fetched from the university's own system at scoring time — not a "
         "self-reported list typed into a profile.",
         icon_name="i-cloud-sync", title_size=12.5, body_size=10)
    card(s, rx, y + rh + 0.16, rwid, rh,
         "Degrades, never breaks",
         "If a university's API is unavailable, the course weight is redistributed "
         "proportionally across the others. The student gets a partial score with "
         "the reason shown — not a zero, and not an error page.",
         icon_name="i-shield", title_size=12.5, body_size=10)

    m2 = len(s.shapes)
    callout(s, M, y + ch + 0.24, CW, 0.72,
            "Job boards match on free text a student writes about themselves. "
            "This matches on what the registrar says they have actually passed.",
            "i-target", "green")

    footer(s, 7)
    set_transition(s, "fade")
    # the bars sweep out of the baseline; everything else fades
    animate(s, waves(s, [m0, m_bars, m_vals, m1, m2]),
            filters=["fade", "wipe(right)", "fade", "fade", "fade"],
            dur=460, gap=150, stagger=110)
    notes(s, """
SPEAKER B  ·  60 sec  ·  140 words

Internship matching is the module with the most direct real-world value — and
the one our university is looking to adopt.

Every posting is scored out of 100. The left panel shows the split: skills 35,
completed courses 25, major 20, interests and location 10 each.

It isn't keyword matching. TF-IDF with cosine similarity, so a term that appears
in every posting counts for very little and a rare skill counts for a lot.

The part we're most pleased with is the top-right box. The course component is
read live through the adapter — the university's record of what the student
actually passed. A job board matches what a student writes about themselves.
We match what the registrar says.

Below it, the honest detail: if the API is down we don't fail. The remaining
weights scale up proportionally and the student sees a partial score with the
reason.
""")
    return s


def slide_rides(prs):
    s = new_slide(prs)
    y = heading(s, "Feature 03  ·  Ride Sharing",
                "Campus-only lifts, tracked live.")

    cards = [
        ("i-car", "Offer and request",
         "A driver-student publishes a route, a departure time and a seat count. "
         "Riders from the same university request a seat; the driver approves or "
         "declines each one.\n\nVehicles are registered to the driver, so a rider "
         "knows what car to expect before it arrives."),
        ("i-map", "Live tracking",
         "Once a ride starts, the driver's position streams over SignalR to "
         "everyone on that ride, on a live map.\n\nA rider waiting at a pickup "
         "point can see where the car actually is, instead of exchanging "
         "messages about it."),
        ("i-lock", "The last seat",
         "Two riders approved at the same instant cannot oversell the car. An "
         "optimistic-concurrency row version means one approval wins and the "
         "other is told cleanly.\n\nWithout it, the driver discovers the problem "
         "in the car park."),
    ]
    cw = (CW - 2 * 0.28) / 3
    marks = []
    for i, (ic, t, b) in enumerate(cards):
        marks.append(len(s.shapes))
        card(s, M + i * (cw + 0.28), y, cw, 2.78, t, b, icon_name=ic,
             title_size=13.5, body_size=10.5)

    # the lifecycle, so the three cards read as one flow rather than three facts
    marks.append(len(s.shapes))
    fy = y + 3.00
    rect(s, M, fy, CW, 0.66, fill=WHITE, line=BORDER, radius=0.10, shadow=True)
    stages = ["Driver offers a route", "Rider requests a seat", "Driver approves",
              "Live tracking en route", "Ride completed"]
    seg = (CW - 0.4) / len(stages)
    for i, stage in enumerate(stages):
        sx = M + 0.2 + i * seg
        tf = tbox(s, sx, fy + 0.08, seg - 0.24, 0.50, anchor="mid")
        para(tf, stage, 10.5, INK_SOFT, bold=(i == 0), font=FONT_SB if i == 0 else FONT,
             first=True, align="c")
        if i < len(stages) - 1:
            tf = tbox(s, sx + seg - 0.26, fy + 0.08, 0.28, 0.50, anchor="mid")
            para(tf, "→", 12, GREEN, bold=True, first=True, align="c")

    marks.append(len(s.shapes))
    callout(s, M, y + 3.86, CW, 0.92,
            "A commercial ride-hailing app can move a student across the city, but it cannot "
            "offer the one guarantee that matters on a campus: every person in that car is a "
            "verified member of your own university, with a real name attached to a real record.",
            "i-shield", "green")

    footer(s, 8)
    set_transition(s, "fade")
    animate(s, waves(s, marks))
    notes(s, """
SPEAKER B  ·  45 sec  ·  110 words

Ride sharing is the most familiar of the three, so I'll be quick on what it does.

A driver publishes a route and a seat count. Riders request, the driver
approves. Once moving, the driver's position streams over SignalR so riders can
watch the car approach.

The engineering point is the third card. If two riders are approved at the same
instant for the last seat, a naive implementation sells it twice. An
optimistic-concurrency row version means one wins and the other gets a clean
message — rather than the driver finding out in the car park.

And the differentiator: Uber can move a student across the city. It can't
guarantee everyone in that car is a verified member of your university.
""")
    return s


def find_shot(*keywords):
    """
    Any screenshot dropped into assets/shots whose filename mentions one of
    these words. Matching on keywords rather than exact names means the file can
    be called whatever came out of the screenshot tool.
    """
    folder = os.path.join(ASSETS, "shots")
    if not os.path.isdir(folder):
        return None
    for name in sorted(os.listdir(folder)):
        low = name.lower()
        if low.endswith((".png", ".jpg", ".jpeg")) and any(k in low for k in keywords):
            return os.path.join(folder, name)
    return None


def feature_slide(prs, number, eyebrow, title, cards, callout_text,
                  callout_icon, shot_keywords, page):
    """
    One of the per-service slides: three cards, a callout, and a screenshot on
    the right if one has been supplied. Without a screenshot the cards simply
    take the full width, so the slide is complete either way.
    """
    s = new_slide(prs)
    y = heading(s, f"Feature {number}  ·  {eyebrow}", title)

    shot = find_shot(*shot_keywords)
    band = 2.45 if shot else 0.0
    left_w = CW - band - (0.20 if shot else 0.0)

    cw = (left_w - 2 * 0.24) / 3
    ch = 2.62
    marks = []
    for i, (ic, t, b) in enumerate(cards):
        marks.append(len(s.shapes))
        card(s, M + i * (cw + 0.24), y, cw, ch, t, b, icon_name=ic,
             title_size=13, body_size=10)

    marks.append(len(s.shapes))
    callout(s, M, y + ch + 0.26, left_w, 0.92, callout_text, callout_icon, "green")

    if shot:
        from PIL import Image
        with Image.open(shot) as im:
            ratio = im.width / im.height
        # portrait shots are phone captures and get the device frame; wide ones
        # are the web portal and get the browser chrome
        framed = os.path.join(ASSETS, "shots", "_framed_" + os.path.basename(shot))
        if ratio < 0.75:
            theme.make_phone_frame(shot, framed)
        else:
            theme.make_browser_frame(shot, framed)
        with Image.open(framed) as im:
            fr = im.width / im.height

        avail_h, avail_w = 3.90, band
        sh = min(avail_h, avail_w / fr)
        sw_ = sh * fr
        pic = picture(s, framed, SW - M - band + (band - sw_) / 2, y, h=sh)
        soft_shadow(pic, blur=24, dist=8, alpha=0.24)

    footer(s, page)
    set_transition(s, "fade")
    animate(s, waves(s, marks))
    return s


def slide_study_groups(prs):
    s = feature_slide(
        prs, "04", "Study Groups",
        "Only your classmates can join.",
        [("i-users", "Scoped to the course",
          "Membership is checked against the university's own enrolment record "
          "through the adapter. Not an invite link, and not a name typed into a "
          "box — the group can only contain people actually taking that course."),
         ("i-bell", "Real-time chat",
          "Messages go over SignalR through the StudyGroupHub, so they appear "
          "instantly for everyone in the group, on the web portal and in the "
          "mobile app alike."),
         ("i-lock", "Capacity that holds",
          "Groups cap at ten members. Two students accepting the last place at "
          "the same instant cannot both get in: an optimistic-concurrency row "
          "version means one succeeds and the other is told cleanly.")],
        "A WhatsApp group cannot tell a classmate from a stranger who was forwarded the "
        "invite link. This can, because it asks the university before letting anyone in.",
        "i-shield", ("group", "study", "chat"), 9)

    notes(s, """
SPEAKER B  ·  45 sec  ·  115 words

Study Groups is the module that shows the enrolment check doing ordinary work.

A student creates a group for a course. Anyone who asks to join is checked
against the university's own enrolment record, live, through the adapter. Not an
invite link — the group can only ever contain people actually taking that course.

Chat is real-time over SignalR, the same hub serving web and mobile.

And the third card is the same concurrency problem as the last seat in a car:
groups cap at ten, and two students accepting the final place at the same moment
cannot both get in.

The callout is the point. A WhatsApp group cannot tell a classmate from a
stranger with the link. This can, because it asks the university.
""")
    return s


def slide_clubs(prs):
    s = feature_slide(
        prs, "05", "Clubs & Organizations",
        "Run a club, not a group chat.",
        [("i-check-badge", "Roles that mean something",
          "President, Vice-President, Officer and Member are real roles with "
          "real permissions, so only officers can post announcements or create "
          "events. A club survives its committee changing."),
         ("i-idea", "Announcements that reach everyone",
          "An officer posts once and every member is notified through the shared "
          "notification service — rather than hoping the message is seen in a "
          "chat thread nobody scrolls back through."),
         ("i-clock", "Events with RSVP",
          "Members respond in advance, so a club knows its numbers before the "
          "day rather than counting whoever turns up at the door.")],
        "Student societies currently run on group chats and spreadsheets, which lose their "
        "history every time the committee changes. Membership, roles and events belong somewhere permanent.",
        "i-idea", ("club", "event", "rsvp"), 10)

    notes(s, """
SPEAKER B  ·  40 sec  ·  100 words

Clubs is the module a student society actually runs on.

The first card is the one that matters. President, Vice-President, Officer and
Member are real roles carrying real permissions — only officers can post an
announcement or create an event. That means a club survives its committee
changing, which a group chat does not.

Announcements go through the same shared notification service every other module
uses. Events carry RSVPs, so a society knows its numbers before the day.

The callout is the honest framing: societies today run on group chats and
spreadsheets, and they lose their history every single year when the committee
hands over.
""")
    return s


def slide_tickets(prs):
    s = feature_slide(
        prs, "06", "Complaints & Ticketing",
        "A complaint with a paper trail.",
        [("i-target", "Routed to a department",
          "A ticket is raised against a category that maps to the department "
          "that owns it, so it reaches the staff who can actually resolve it "
          "instead of a shared inbox nobody is accountable for."),
         ("i-clock", "Tracked to resolution",
          "Open, In Progress, Waiting for Student, Resolved, Closed or Rejected. "
          "The student can see the state of their own complaint at any moment, "
          "and every response is kept."),
         ("i-file", "Evidence and an audit trail",
          "Attachments go with the ticket, and every administrative action is "
          "written to the audit log — so what happened, and who did it, is "
          "answerable after the fact.")],
        "An email to a department has no status, no owner and no record. A student chasing "
        "a complaint should never have to ask whether anyone received it.",
        "i-check-badge", ("ticket", "complaint", "support"), 11)

    notes(s, """
SPEAKER B  ·  40 sec  ·  105 words

Complaints and Ticketing is the least glamorous module and the one a university
administrator asks about first.

A student raises a ticket against a category, and that category maps to the
department that owns it — so it reaches staff who can actually resolve it, not a
shared inbox where nothing is anyone's job.

It then has a real lifecycle: Open, In Progress, Waiting for Student, Resolved,
Closed, Rejected. The student can see that state at any time, and every response
is kept with the ticket.

Attachments carry evidence, and administrative actions are written to the audit
log, so what happened and who did it is answerable afterwards.
""")
    return s


def slide_security(prs):
    s = new_slide(prs)
    y = heading(s, "Security & data protection",
                "What we enforced, and what we are honest about.",
                sub="Smart Attendance records where a named student physically was. "
                    "That raised the bar for what we had to justify.")

    lw = (CW - 0.30) / 2
    ch = 4.38
    m0 = len(s.shapes)

    rect(s, M, y, lw, ch, fill=WHITE, line=BORDER, radius=0.14, shadow=True)
    rect(s, M, y, lw, 0.06, fill=GREEN)
    icon(s, "i-shield", "16a34a", M + 0.30, y + 0.28, 0.30)
    tf = tbox(s, M + 0.72, y + 0.26, lw - 1.0, 0.34, anchor="mid")
    para(tf, "Enforced in the system", 14, INK, bold=True, font=FONT_SB, first=True)
    tf = tbox(s, M + 0.30, y + 0.76, lw - 0.60, ch - 1.0)
    enforced = [
        ("Tenant isolation in the database. ", "Every service row carries a UniversityCode, so isolation survives an application bug."),
        ("A read-only adapter. ", "UniConnect cannot write to university records, by construction."),
        ("Two-factor authentication (web). ", "TOTP under RFC 6238, enrolled by QR so the secret is never typed; the mobile API refuses 2FA accounts."),
        ("Two schemes, one identity. ", "Cookie for web, JWT for mobile, one Identity store."),
        ("Role-based authorization, ", "with denied attempts written to the audit log."),
        ("Anti-forgery tokens ", "on every state-changing web action."),
        ("HTTPS redirection and HSTS, ", "plus ASP.NET Identity password hashing."),
        ("Security-stamp revalidation. ", "A disabled account loses access immediately."),
        ("Uploads stored under non-guessable names.", ""),
    ]
    for i, (lead, rest) in enumerate(enforced):
        bullet(tf, rest, 10.5, MUTED, first=(i == 0), before=0 if i == 0 else 7, bold_lead=lead)

    m1 = len(s.shapes)
    x2 = M + lw + 0.30
    rect(s, x2, y, lw, ch, fill=WHITE, line=BORDER, radius=0.14, shadow=True)
    rect(s, x2, y, lw, 0.06, fill=AMBER)
    icon(s, "i-alert", "d97706", x2 + 0.30, y + 0.28, 0.30)
    tf = tbox(s, x2 + 0.72, y + 0.26, lw - 1.0, 0.34, anchor="mid")
    para(tf, "Known gaps — documented, not hidden", 14, INK, bold=True,
         font=FONT_SB, first=True)
    tf = tbox(s, x2 + 0.30, y + 0.76, lw - 0.60, ch - 1.0)
    gaps = [
        ("No retention or erasure policy. ", "Attendance records and GPS coordinates persist indefinitely, with no self-service export or delete."),
        ("Uneven disclosure. ", "Both clients tell the student their location is checked; the device fingerprint is collected without being surfaced."),
        ("No formal DPIA ", "was carried out."),
        ("Files on local disk, ", "rather than object storage with per-object access control."),
        ("Authorization is evidenced by code review, ", "not yet by a formal test matrix."),
        ("No load testing, ", "so graceful degradation is demonstrated by construction rather than measurement."),
    ]
    for i, (lead, rest) in enumerate(gaps):
        bullet(tf, rest, 10.5, MUTED, first=(i == 0), before=0 if i == 0 else 9, bold_lead=lead)

    footer(s, 12)
    set_transition(s, "fade")
    animate(s, waves(s, [m0, m1]))
    notes(s, """
SPEAKER B  ·  60 sec  ·  140 words

We'll be direct here, because attendance changes what kind of system this is.
Recording GPS at check-in means storing where a named student physically was.

On the left, what's enforced. The first matters most: tenant isolation lives in
the database. Every service row carries a university code, so one university's
staff can't reach another's data even if we made a mistake in application code.

Second, the adapter is read-only by construction — no code path writes to a
university's record, so a bug here cannot corrupt one.

Third, two-factor authentication. TOTP, and the enrolment detail matters:
the student scans a QR code, so the shared secret is never typed by hand.

On the right, what we haven't done. No retention or erasure policy. We tell
students we check location, not that we fingerprint the device. No formal DPIA.

We could have left this slide out. A system that records student location should
be presented with its gaps stated.
""")
    return s


def slide_comparison(prs):
    s = new_slide(prs)
    y = heading(s, "Differentiation",
                "No single existing product covers this ground.")

    cols = ["LMS\nCanvas, Blackboard", "Messaging\nWhatsApp, Discord",
            "Ride-hailing\napps", "Job boards\nLinkedIn, Indeed", "UniConnect"]
    rows = [
        ("Verified course enrolment behind every action", ["y", "n", "n", "n", "y"]),
        ("Attendance verified on more than one factor", ["p", "n", "n", "n", "y"]),
        ("Trust scoped to your own university", ["y", "p", "n", "n", "y"]),
        ("Matching on actual completed coursework", ["n", "n", "n", "p", "y"]),
        ("Runs above the university, owning no academic data", ["n", "n", "n", "n", "y"]),
        ("All of it behind one verified identity", ["n", "n", "n", "n", "y"]),
    ]

    label_w = 4.35
    col_w = (CW - label_w) / len(cols)
    hdr_h = 0.66
    row_h = 0.50
    marks = [len(s.shapes)]

    # header
    for i, c in enumerate(cols):
        cx = M + label_w + i * col_w
        last = (i == len(cols) - 1)
        if last:
            rect(s, cx, y - 0.06, col_w, hdr_h + 0.06 + len(rows) * row_h + 0.10,
                 fill=GREEN_TINT, line=GREEN_SOFT, radius=0.10)
        tf = tbox(s, cx + 0.06, y, col_w - 0.12, hdr_h, anchor="mid")
        para(tf, c, 9.5 if not last else 11.5,
             GREEN_DEEP if last else MUTED, bold=last, font=FONT_SB if last else FONT,
             first=True, align="c", spacing=1.08)

    ty = y + hdr_h
    line_h(s, M, ty, CW, BORDER, 1.0)
    marks.append(len(s.shapes))

    for r, (label, cells) in enumerate(rows):
        ry = ty + r * row_h
        if r % 2 == 0:
            rect(s, M, ry, label_w, row_h, fill=WHITE)
        tf = tbox(s, M + 0.10, ry, label_w - 0.24, row_h, anchor="mid")
        para(tf, label, 10.5, INK_SOFT, first=True)
        for i, mk in enumerate(cells):
            cx = M + label_w + i * col_w
            mid = cx + col_w / 2
            if mk == "y":
                icon(s, "i-check-circle", "16a34a", mid - 0.115, ry + row_h / 2 - 0.115, 0.23)
            elif mk == "p":
                tf = tbox(s, cx, ry, col_w, row_h, anchor="mid")
                para(tf, "partial", 9, AMBER, bold=True, font=FONT_SB, first=True, align="c")
            else:
                tf = tbox(s, cx, ry, col_w, row_h, anchor="mid")
                para(tf, "—", 12, "cbd5e1", first=True, align="c")
        line_h(s, M, ry + row_h, CW, BORDER_SFT, 0.75)

    marks.append(len(s.shapes))
    ly = ty + len(rows) * row_h + 0.30
    callout(s, M, ly, CW, 0.80,
            "Each column does its own job well. The gap is that a student's enrolment — the one "
            "fact that makes any of these trustworthy inside a university — is visible to none of them.",
            "i-idea", "green")

    footer(s, 13)
    set_transition(s, "fade")
    animate(s, waves(s, marks))
    notes(s, """
SPEAKER A  ·  50 sec  ·  115 words

We're not claiming these products are bad. Each does its own job better than we
do. Canvas is a better LMS than anything we built; LinkedIn has a vastly larger
job pool.

The claim is coverage. Read across the rows.

An LMS knows enrolment, but its attendance is usually a single code — partial.
WhatsApp has no idea who's enrolled; anyone with the link joins. A ride app has
no campus context. LinkedIn matches, but on text the student wrote about
themselves.

And the bottom two rows are empty until the last column. Nothing there sits
above the university's own system without absorbing its data, and nothing puts
all of this behind one verified identity.
""")
    return s


def slide_evidence(prs):
    s = new_slide(prs)
    y = heading(s, "Evidence", "What was actually built and verified.")

    stats = [
        ("290", "passing tests"),
        ("~32,300", "lines of hand-written C#"),
        ("37", "database tables"),
        ("174", "controller actions"),
    ]
    stats2 = [
        ("22", "service classes"),
        ("6", "SignalR hubs"),
        ("33", "EF Core migrations"),
        ("2", "provider implementations"),
    ]
    sw_ = (CW - 3 * 0.26) / 4
    m0 = len(s.shapes)
    for i, (v, l) in enumerate(stats):
        stat(s, M + i * (sw_ + 0.26), y, sw_, 1.32, v, l)
    m1 = len(s.shapes)
    for i, (v, l) in enumerate(stats2):
        stat(s, M + i * (sw_ + 0.26), y + 1.50, sw_, 1.32, v, l)

    m2 = len(s.shapes)
    by = y + 3.10
    rect(s, M, by, CW, 1.20, fill=GREEN_DEEP, radius=0.12, shadow=True)
    icon(s, "i-check-badge", "ffffff", M + 0.34, by + 0.42, 0.36)
    tf = tbox(s, M + 0.92, by + 0.20, CW - 1.3, 0.82, anchor="mid")
    para(tf, "Tested against a real relational database, not an in-memory stub.",
         13, WHITE, bold=True, font=FONT_SB, first=True)
    para(tf, "EF Core's in-memory provider does not enforce unique indexes or foreign keys — so a suite built on it "
             "passes silently for exactly the defects most likely to reach production: a duplicate check-in, a second "
             "application to the same internship, an orphaned membership row.",
         10.5, "b7ecc9", before=4, spacing=1.16)

    footer(s, 14)
    set_transition(s, "fade")
    animate(s, waves(s, [m0, m1, m2]))
    notes(s, """
SPEAKER A  ·  45 sec  ·  110 words

Briefly, the scale — a defence should point at something concrete.

About twenty-eight thousand seven hundred lines of hand-written C#, excluding
generated migrations. Thirty-seven tables, a hundred and seventy-four controller
actions, six SignalR hubs. Two provider implementations — the adapter claim in
numeric form.

Two hundred and seventy-one passing tests. The band at the bottom is what we'd
like to be asked about.

Most student suites use EF Core's in-memory provider — fast, easy, and it
doesn't enforce unique indexes or foreign keys. So it passes for exactly the
bugs most likely to reach production. We run against a real database. Slower.
It tests something real.
""")
    return s


def slide_future(prs):
    s = new_slide(prs)
    y = heading(s, "Future work", "The next steps, in the order we would take them.")

    items = [
        ("i-bell", "The mobile live layer",
         "SignalR-backed chat, live ride tracking and push notifications, bringing "
         "the mobile client up to the real-time behaviour the web portal already has."),
        ("i-car", "Full mobile coverage",
         "Clubs and Complaints as complete mobile modules, plus a dedicated "
         "endpoint for the cross-module My Activity view, and code entry so "
         "two-factor accounts can sign in on the app."),
        ("i-cloud-sync", "Close the partner-API gaps",
         "An instructor and staff directory, and either a bulk student endpoint or "
         "a documented workaround. Both are already raised with the partner university."),
        ("i-check-circle", "Test depth",
         "Extend end-to-end HTTP testing beyond authentication to every module, "
         "and add a formal authorization matrix asserting each role reaches "
         "exactly the endpoints it should."),
        ("i-shield", "Data protection",
         "A retention and erasure policy, self-service export and delete, full "
         "disclosure at the point of collection, and a formal DPIA."),
        ("i-users", "Usability evaluation",
         "No study was run with real students, and the system has not been load "
         "tested. Both would turn design decisions we argue for into results we "
         "can measure."),
    ]
    cw = (CW - 2 * 0.28) / 3
    ch = 1.95
    marks = []
    for i, (ic, t, b) in enumerate(items):
        if i % 3 == 0:
            marks.append(len(s.shapes))
        cx = M + (i % 3) * (cw + 0.28)
        cy = y + (i // 3) * (ch + 0.26)
        card(s, cx, cy, cw, ch, t, b, icon_name=ic, title_size=12.5, body_size=10)

    footer(s, 15)
    set_transition(s, "fade")
    animate(s, waves(s, marks))
    notes(s, """
SPEAKER B  ·  40 sec  ·  105 words  ·  Don't read all six — lead with mobile

Six items; I'll lead with the one we'd do next.

The mobile live layer. The web portal already has real-time chat, live ride
tracking and instant notifications over SignalR. The mobile client doesn't yet,
and closing that gap is the single biggest improvement to how the app feels.

After that: finish mobile coverage for rides, clubs and complaints; close the
two capabilities genuinely missing from the partner university's API, both
already raised with them; deepen the test suite with end-to-end tests and a
formal authorization matrix; and complete the data-protection work from the
previous slide.

The last one is worth naming. We argue for our interface decisions from
convention. We have not measured them against real students.
""")
    return s


def slide_close(prs):
    s = new_slide(prs, dark=True)

    # Centred, unlike every other slide. The change of axis is what makes it
    # read as an ending rather than as one more content slide.
    n0 = len(s.shapes)
    mid = SW / 2
    glow = os.path.join(ASSETS, "glow.png")
    if os.path.exists(glow):
        picture(s, glow, mid - 1.55, 0.42, w=3.10, h=3.10)
    rect(s, mid - 0.44, 1.28, 0.88, 0.88, fill=GREEN, radius=0.23)
    icon(s, "i-graduation", "ffffff", mid - 0.23, 1.49, 0.46)
    n1 = len(s.shapes)

    tf = tbox(s, M, 2.42, CW, 1.10)
    para(tf, "Thank you.", 60, WHITE, bold=True, font=FONT_SB, first=True, align="c")

    tf = tbox(s, M, 3.56, CW, 0.42)
    para(tf, "We are happy to take questions.", 18.5, GREEN_SOFT,
         first=True, align="c")
    n2 = len(s.shapes)

    # the headline numbers, at the moment the room is paying most attention
    stats = [("9", "services"), ("290", "passing tests"),
             ("2", "universities integrated"), ("~32,300", "lines of C#")]
    band_w, band_h = 9.40, 1.16
    bx = mid - band_w / 2
    # translucent, so the constellation reads through it as glass
    rect(s, bx, 4.20, band_w, band_h, fill="0b3d21", radius=0.16,
         line="2f7a4f", lw=0.75, alpha=0.42)
    seg = band_w / len(stats)
    for i, (value, label) in enumerate(stats):
        cx = bx + i * seg
        tf = tbox(s, cx, 4.34, seg, 0.52, anchor="mid")
        para(tf, value, 26, WHITE, bold=True, font=FONT_SB, first=True, align="c")
        tf = tbox(s, cx, 4.92, seg, 0.32, anchor="mid")
        para(tf, label, 10.5, "9fdcb6", first=True, align="c")
        if i:
            rect(s, cx, 4.44, 0.012, 0.68, fill="3d8a5e", alpha=0.55)
    n3 = len(s.shapes)

    line_h(s, mid - 0.95, 5.78, 1.9, GREEN_LT, 2.6)

    tf = tbox(s, M, 6.04, CW, 0.80)
    para(tf, "Mohamad Sabbagh  ·  1123206          Mohamad Ali Rida  ·  1122824",
         13, WHITE, bold=True, font=FONT_SB, first=True, align="c")
    para(tf, "Supervisor:  Zahraa Sleem      ·      UniConnect  —  Connect. Study. Commute.",
         11, "8fd0a8", align="c", before=6)

    set_transition(s, "fade")
    animate(s, waves(s, [n0, n1, n2, n3]), dur=520, gap=210)

    notes(s, """
BOTH  ·  10 sec  ·  10 words

Thank you — we're happy to take questions.

(Everything below is reference for the Q&A. Not spoken.)

--------------------------------------------------------------------
LIKELY QUESTIONS, AND THE SHORT ANSWERS
--------------------------------------------------------------------

"Your report lists two-factor authentication as future work."
  Correct at the time of writing. It was finished after the report was
  submitted, which is why it appears on the security slide rather than the
  future-work one. Offer to demonstrate the QR enrolment.

"Does two-factor work on the mobile app?"
  No, and that is deliberate rather than unfinished. The mobile API validates
  the password only, so leaving it alone would have meant a student could
  enable 2FA on the web and still be let into the app with a password alone —
  the second factor would have been decorative. The API refuses those accounts
  with a clear message instead, and the enrolment page says so before the
  student opts in. Accepting a code in the login request is the next step; it
  needs a mobile release.

"How do you know an authenticator app will actually accept your QR code?"
  We compute the TOTP independently — in the test suite and again in an
  end-to-end script — from nothing but the shared key and the clock, which is
  exactly what a phone has. The server accepts those codes. We also loaded the
  QR library outside the browser and had it encode the real provisioning URI.

"What if a student loses their phone?"
  Ten single-use recovery codes, shown once behind an explicit confirmation,
  with copy, download and print. If those are gone too, an administrator can
  reset the second factor — and that reset writes an audit entry naming the
  administrator who did it, which is what keeps the escape hatch from being a
  back door.

"Why not just build this into the LMS?"
  Because it would require the university to let a third party write into its
  system of record. Our whole architecture exists to avoid asking for that.

"What stops a student spoofing GPS?"
  Nothing at the OS level — a rooted device can lie. That is exactly why GPS is
  one of five checks rather than the only one, and why device-fingerprint reuse
  is checked alongside it. We would not claim it is unspoofable.

"How long would onboarding a new university take?"
  Implementing one interface. For the partner university it was one class. The
  real cost is not the code, it is getting API access agreed.

"Is the data safe?"
  Isolation is enforced at the database level, and the adapter cannot write to
  university systems. But we listed real gaps on the security slide — no
  retention policy, no DPIA — and we would not deploy without closing them.

"Why is the mobile app incomplete?"
  We lost two months of work when access to a development account was lost. We
  chose to finish three modules properly rather than start six badly, and
  documented the rest as future work.

"What was the hardest part?"
  Integrating with a real external API we did not control. It had different
  field names, nested course data, and two endpoints that simply do not exist.
  That is also what proved the adapter was worth building.
""")
    return s


# =================================================================== main ===
def prepare_assets():
    os.makedirs(ASSETS, exist_ok=True)
    make_backgrounds()

    # icons come from the sprite renderer that lives beside this file
    renderer = os.path.join(HERE, "render_icons.py")
    subprocess.run([sys.executable, renderer], check=True)

    # the attendance screenshot is pulled out of the report so the deck and the
    # report cannot drift apart
    # Screenshots come out of the report, so the deck and the report can never
    # show different builds of the app.
    #   image7  Figure 5 — web portal, Home dashboard
    #   image8  Figure 6 — mobile app, Home tab
    #   image10 Figure 8 — mobile app, Attendance check-in
    import zipfile
    wanted = {"image7.png": "web_home.png",
              "image8.png": "mobile_home.png",
              "image10.png": "mobile_checkin.png"}
    missing = [d for d in wanted.values() if not os.path.exists(os.path.join(ASSETS, d))]
    if missing:
        docx = os.path.join(ROOT, "UniConnect_Final_Report.docx")
        if os.path.exists(docx):
            with zipfile.ZipFile(docx) as z:
                for src_name, dst_name in wanted.items():
                    entry = f"word/media/{src_name}"
                    if entry in z.namelist():
                        with z.open(entry) as src, \
                             open(os.path.join(ASSETS, dst_name), "wb") as dst:
                            dst.write(src.read())
            print(f"extracted {len(wanted)} screenshots from the report")

    def framed(name, fn, out):
        p = os.path.join(ASSETS, name)
        if os.path.exists(p):
            fn(p, os.path.join(ASSETS, out))

    framed("mobile_checkin.png", theme.make_phone_frame, "mobile_checkin_framed.png")
    framed("mobile_home.png", theme.make_phone_frame, "mobile_home_framed.png")
    framed("web_home.png", theme.make_browser_frame, "web_home_framed.png")
    theme.make_glow(os.path.join(ASSETS, "glow.png"))


def build():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_architecture(prs)
    slide_stack(prs)
    slide_attendance(prs)
    slide_internships(prs)
    slide_rides(prs)
    slide_study_groups(prs)
    slide_clubs(prs)
    slide_tickets(prs)
    slide_security(prs)
    slide_comparison(prs)
    slide_evidence(prs)
    slide_future(prs)
    slide_close(prs)

    prs.save(OUT)
    print(f"saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    prepare_assets()
    build()
